#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Normalisation de fichiers hétérogènes pour ingestion par une IA (ex. NotebookLM).

Entrée  : dossier ./sources_brutes (par défaut)
Sortie  : dossier ./sources_normalisees (par défaut)

Règle générale : 1 fichier en entrée -> 1 fichier en sortie (.txt ou .csv UTF-8).

Règles de conversion :
  .xlsx / .xlsm / .xls  -> 1 CSV si une feuille, sinon 1 .txt (feuilles en CSV délimitées)
  .csv  / .tsv          -> CSV ré-encodé UTF-8, séparateur détecté puis normalisé en virgule
  .pdf                  -> fichier .txt (texte extrait page par page)
  .docx                 -> fichier .txt (paragraphes + tableaux convertis en lignes pipe)
  .pptx                 -> fichier .txt (texte des diapositives)
  .html / .htm          -> fichier .txt (texte sans balises)
  .json / .xml / .yaml  -> fichier .txt ré-encodé UTF-8 (contenu inchangé)
  .sql / .txt / .md / code (.py, .js, ...) -> fichier .txt ré-encodé UTF-8
  .zip                  -> 1 seul .txt consolidé (fichiers internes délimités)

Chaque fichier de sortie commence par un en-tête de métadonnées non ambigu.
Un fichier _MANIFEST.csv récapitule toutes les conversions.
"""

from __future__ import annotations

import argparse
import csv
import io
import re
import sys
import unicodedata
from datetime import datetime
from pathlib import Path

# ---------------------------------------------------------------------------
# Configuration des types de fichiers
# ---------------------------------------------------------------------------

EXT_ZIP = {".zip"}
EXT_TABLEUR = {".xlsx", ".xlsm", ".xls"}
EXT_CSV = {".csv", ".tsv"}
EXT_PDF = {".pdf"}
EXT_DOCX = {".docx"}
EXT_PPTX = {".pptx"}
EXT_HTML = {".html", ".htm"}
EXT_TEXTE = {
    ".txt", ".md", ".sql", ".json", ".xml", ".yaml", ".yml", ".ini", ".cfg",
    ".log", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".cs", ".sh", ".ps1",
    ".bat", ".r", ".rb", ".php", ".go", ".rs", ".toml", ".env", ".csvlog",
    ".qvs",  # scripts Qlik
}

# ---------------------------------------------------------------------------
# Utilitaires
# ---------------------------------------------------------------------------


def slugifier(nom: str) -> str:
    """Transforme un nom en identifiant ASCII sûr (minuscules, underscores)."""
    nom = unicodedata.normalize("NFKD", nom).encode("ascii", "ignore").decode("ascii")
    nom = re.sub(r"[^\w\-]+", "_", nom).strip("_").lower()
    return nom or "sans_nom"


def lire_texte_brut(chemin: Path) -> str:
    """Lit un fichier texte en détectant l'encodage."""
    donnees = chemin.read_bytes()
    try:
        return donnees.decode("utf-8-sig")
    except UnicodeDecodeError:
        pass
    try:
        from charset_normalizer import from_bytes

        resultat = from_bytes(donnees).best()
        if resultat is not None:
            return str(resultat)
    except ImportError:
        pass
    return donnees.decode("latin-1", errors="replace")


def entete_metadonnees(source: Path, type_doc: str, complement: str = "") -> str:
    lignes = [
        "=" * 78,
        "DOCUMENT NORMALISE POUR ANALYSE PAR IA",
        f"Fichier source : {source.name}",
        f"Type d'origine : {type_doc}",
        f"Date de conversion : {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
    ]
    if complement:
        lignes.append(complement)
    lignes.append("=" * 78)
    return "\n".join(lignes) + "\n\n"


def ecrire_txt(destination: Path, contenu: str) -> None:
    destination.write_text(contenu, encoding="utf-8", newline="\n")


# ---------------------------------------------------------------------------
# Convertisseurs
# ---------------------------------------------------------------------------


def convertir_tableur(source: Path, dossier_sortie: Path) -> list[Path]:
    """Excel -> un seul fichier de sortie.

    - Classeur à une feuille : un CSV UTF-8.
    - Classeur à plusieurs feuilles : un seul .txt contenant chaque feuille
      au format CSV, délimitée par une ligne '>>> FEUILLE : <nom>'.
    """
    import pandas as pd

    feuilles = pd.read_excel(source, sheet_name=None, dtype=str)
    base = slugifier(source.stem)

    if len(feuilles) == 1:
        df = next(iter(feuilles.values()))
        df = df.dropna(how="all").dropna(axis=1, how="all")
        destination = dossier_sortie / f"{base}.csv"
        df.to_csv(destination, index=False, encoding="utf-8", lineterminator="\n")
        return [destination]

    morceaux = [entete_metadonnees(
        source, "Classeur Excel",
        f"Nombre de feuilles : {len(feuilles)} — chaque feuille est au format CSV,\n"
        "délimitée par une ligne '>>> FEUILLE : <nom>'.")]
    for nom_feuille, df in feuilles.items():
        df = df.dropna(how="all").dropna(axis=1, how="all")
        morceaux.append(f"\n>>> FEUILLE : {nom_feuille}\n")
        morceaux.append(df.to_csv(index=False, lineterminator="\n").rstrip())
    destination = dossier_sortie / f"{base}.txt"
    ecrire_txt(destination, "\n".join(morceaux) + "\n")
    return [destination]


def convertir_csv(source: Path, dossier_sortie: Path) -> list[Path]:
    """CSV/TSV -> CSV UTF-8 avec séparateur virgule."""
    texte = lire_texte_brut(source)
    extrait = texte[:8192]
    try:
        dialecte = csv.Sniffer().sniff(extrait, delimiters=",;\t|")
        separateur = dialecte.delimiter
    except csv.Error:
        separateur = "\t" if source.suffix.lower() == ".tsv" else ","

    lignes = list(csv.reader(io.StringIO(texte), delimiter=separateur))
    destination = dossier_sortie / f"{slugifier(source.stem)}.csv"
    with destination.open("w", encoding="utf-8", newline="") as f:
        csv.writer(f, lineterminator="\n").writerows(lignes)
    return [destination]


def convertir_pdf(source: Path, dossier_sortie: Path) -> list[Path]:
    """PDF -> texte, page par page."""
    import pdfplumber

    morceaux: list[str] = []
    with pdfplumber.open(source) as pdf:
        nb_pages = len(pdf.pages)
        for i, page in enumerate(pdf.pages, start=1):
            texte = page.extract_text() or "[Page sans texte extractible]"
            morceaux.append(f"--- PAGE {i} / {nb_pages} ---\n{texte}")
    contenu = entete_metadonnees(source, "PDF", f"Nombre de pages : {nb_pages}")
    contenu += "\n\n".join(morceaux) + "\n"
    destination = dossier_sortie / f"{slugifier(source.stem)}.txt"
    ecrire_txt(destination, contenu)
    return [destination]


def convertir_docx(source: Path, dossier_sortie: Path) -> list[Path]:
    """DOCX -> texte (paragraphes + tableaux en lignes pipe)."""
    import docx

    document = docx.Document(str(source))
    morceaux: list[str] = []

    for para in document.paragraphs:
        texte = para.text.strip()
        if not texte:
            continue
        style = (para.style.name or "").lower()
        if style.startswith("heading") or style.startswith("titre"):
            niveau = re.search(r"\d+", style)
            prefixe = "#" * (int(niveau.group()) if niveau else 1)
            morceaux.append(f"\n{prefixe} {texte}")
        else:
            morceaux.append(texte)

    for idx, table in enumerate(document.tables, start=1):
        lignes = [f"\n--- TABLEAU {idx} ---"]
        for rangee in table.rows:
            cellules = [c.text.strip().replace("\n", " ") for c in rangee.cells]
            lignes.append(" | ".join(cellules))
        morceaux.append("\n".join(lignes))

    contenu = entete_metadonnees(source, "Document Word (.docx)")
    contenu += "\n".join(morceaux) + "\n"
    destination = dossier_sortie / f"{slugifier(source.stem)}.txt"
    ecrire_txt(destination, contenu)
    return [destination]


def convertir_pptx(source: Path, dossier_sortie: Path) -> list[Path]:
    """PPTX -> texte par diapositive."""
    from pptx import Presentation

    presentation = Presentation(str(source))
    morceaux: list[str] = []
    for i, diapo in enumerate(presentation.slides, start=1):
        textes = []
        for forme in diapo.shapes:
            if forme.has_text_frame:
                t = forme.text_frame.text.strip()
                if t:
                    textes.append(t)
        morceaux.append(f"--- DIAPOSITIVE {i} ---\n" + ("\n".join(textes) or "[vide]"))
    contenu = entete_metadonnees(
        source, "Présentation PowerPoint (.pptx)",
        f"Nombre de diapositives : {len(presentation.slides)}",
    )
    contenu += "\n\n".join(morceaux) + "\n"
    destination = dossier_sortie / f"{slugifier(source.stem)}.txt"
    ecrire_txt(destination, contenu)
    return [destination]


def convertir_html(source: Path, dossier_sortie: Path) -> list[Path]:
    """HTML -> texte sans balises."""
    from bs4 import BeautifulSoup

    soupe = BeautifulSoup(lire_texte_brut(source), "html.parser")
    for balise in soupe(["script", "style", "noscript"]):
        balise.decompose()
    texte = re.sub(r"\n{3,}", "\n\n", soupe.get_text(separator="\n")).strip()
    contenu = entete_metadonnees(source, "Page HTML") + texte + "\n"
    destination = dossier_sortie / f"{slugifier(source.stem)}.txt"
    ecrire_txt(destination, contenu)
    return [destination]


def convertir_texte(source: Path, dossier_sortie: Path) -> list[Path]:
    """Fichier texte (sql, md, json, code...) -> .txt ré-encodé UTF-8."""
    extension = source.suffix.lower().lstrip(".") or "texte"
    contenu = entete_metadonnees(source, f"Fichier texte (.{extension})")
    contenu += lire_texte_brut(source).rstrip() + "\n"
    nom = slugifier(source.stem)
    if extension not in ("txt",):
        nom = f"{nom}__{extension}"
    destination = dossier_sortie / f"{nom}.txt"
    ecrire_txt(destination, contenu)
    return [destination]


def convertir_zip(source: Path, dossier_sortie: Path) -> list[Path]:
    """ZIP -> un seul .txt consolidé contenant tous les fichiers texte de l'archive.

    Chaque fichier interne est délimité par une ligne '>>> FICHIER : <chemin>'.
    Les fichiers binaires (images, etc.) sont ignorés.
    """
    import zipfile

    fichiers: list[tuple[str, str]] = []
    ignores = 0

    with zipfile.ZipFile(source) as archive:
        for info in sorted(archive.infolist(), key=lambda i: i.filename):
            if info.is_dir():
                continue
            extension = Path(info.filename).suffix.lower()
            if extension not in EXT_TEXTE and extension not in EXT_CSV:
                ignores += 1
                continue
            donnees = archive.read(info)
            try:
                texte = donnees.decode("utf-8-sig")
            except UnicodeDecodeError:
                texte = donnees.decode("latin-1", errors="replace")
            fichiers.append((info.filename, texte))

    morceaux = [entete_metadonnees(
        source, "Archive ZIP",
        f"Contenu : {len(fichiers)} fichier(s) texte consolidé(s) — chaque fichier\n"
        "interne est délimité par une ligne '>>> FICHIER : <chemin dans l'archive>'.")]
    for chemin_interne, texte in fichiers:
        morceaux.append(f"\n>>> FICHIER : {chemin_interne}\n")
        morceaux.append(texte.rstrip())

    destination = dossier_sortie / f"{slugifier(source.stem)}.txt"
    ecrire_txt(destination, "\n".join(morceaux) + "\n")
    if ignores:
        print(f"           ({ignores} fichier(s) binaire(s) ignoré(s) dans l'archive)")
    return [destination]


# ---------------------------------------------------------------------------
# Découpage des fichiers trop volumineux
# ---------------------------------------------------------------------------

# NotebookLM accepte ~500 000 mots par source ; on garde une marge de sécurité.
MAX_MOTS_DEFAUT = 400_000


def decouper_fichier(chemin: Path, max_mots: int) -> list[Path]:
    """Découpe un fichier de sortie en plusieurs parties si trop volumineux.

    La coupe se fait sur des limites de lignes. Pour un CSV, la ligne d'en-tête
    est répétée au début de chaque partie. Pour un .txt, chaque partie commence
    par un bandeau 'PARTIE x / y'. Retourne la liste des fichiers finaux.
    """
    if max_mots <= 0:
        return [chemin]

    lignes = chemin.read_text(encoding="utf-8").splitlines()
    nb_mots_total = sum(len(ligne.split()) for ligne in lignes)
    if nb_mots_total <= max_mots:
        return [chemin]

    est_csv = chemin.suffix.lower() == ".csv"
    entete_csv = lignes[0] if est_csv and lignes else None
    corps = lignes[1:] if entete_csv is not None else lignes

    # Répartition du corps en blocs de max_mots (coupes sur lignes entières)
    blocs: list[list[str]] = []
    bloc_courant: list[str] = []
    mots_bloc = 0
    for ligne in corps:
        mots_ligne = len(ligne.split())
        if bloc_courant and mots_bloc + mots_ligne > max_mots:
            blocs.append(bloc_courant)
            bloc_courant = []
            mots_bloc = 0
        bloc_courant.append(ligne)
        mots_bloc += mots_ligne
    if bloc_courant:
        blocs.append(bloc_courant)

    nb_parties = len(blocs)
    sorties: list[Path] = []
    for i, bloc in enumerate(blocs, start=1):
        destination = chemin.with_name(
            f"{chemin.stem}_partie_{i}_sur_{nb_parties}{chemin.suffix}")
        morceaux: list[str] = []
        if est_csv and entete_csv is not None:
            morceaux.append(entete_csv)
        else:
            morceaux.extend([
                "=" * 78,
                f"PARTIE {i} / {nb_parties} du document '{chemin.name}'",
                "Document découpé car trop volumineux pour une source unique.",
                "=" * 78,
                "",
            ])
        morceaux.extend(bloc)
        destination.write_text("\n".join(morceaux) + "\n",
                               encoding="utf-8", newline="\n")
        sorties.append(destination)

    chemin.unlink()
    print(f"           (fichier volumineux : {nb_mots_total:,} mots, "
          f"découpé en {nb_parties} parties)".replace(",", " "))
    return sorties


# ---------------------------------------------------------------------------
# Orchestration
# ---------------------------------------------------------------------------


def choisir_convertisseur(extension: str):
    if extension in EXT_TABLEUR:
        return convertir_tableur, "tableur"
    if extension in EXT_CSV:
        return convertir_csv, "csv"
    if extension in EXT_PDF:
        return convertir_pdf, "pdf"
    if extension in EXT_DOCX:
        return convertir_docx, "docx"
    if extension in EXT_PPTX:
        return convertir_pptx, "pptx"
    if extension in EXT_HTML:
        return convertir_html, "html"
    if extension in EXT_TEXTE:
        return convertir_texte, "texte"
    if extension in EXT_ZIP:
        return convertir_zip, "zip"
    return None, None


def normaliser_dossier(dossier_entree: Path, dossier_sortie: Path,
                       max_mots: int = MAX_MOTS_DEFAUT) -> int:
    if not dossier_entree.is_dir():
        print(f"ERREUR : le dossier d'entrée n'existe pas : {dossier_entree}")
        return 1

    dossier_sortie.mkdir(parents=True, exist_ok=True)
    manifeste: list[dict[str, str]] = []
    nb_ok = nb_ignores = nb_erreurs = 0

    fichiers = sorted(p for p in dossier_entree.rglob("*") if p.is_file())
    if not fichiers:
        print(f"Aucun fichier trouvé dans {dossier_entree}.")
        return 0

    print(f"Normalisation de {len(fichiers)} fichier(s) : "
          f"{dossier_entree} -> {dossier_sortie}\n")

    for fichier in fichiers:
        extension = fichier.suffix.lower()
        convertisseur, categorie = choisir_convertisseur(extension)
        relatif = fichier.relative_to(dossier_entree)

        if convertisseur is None:
            print(f"  [IGNORÉ ] {relatif} (extension non gérée : {extension or 'aucune'})")
            manifeste.append({
                "fichier_source": str(relatif), "categorie": "non_geree",
                "statut": "ignore", "fichiers_sortie": "", "erreur": "",
            })
            nb_ignores += 1
            continue

        try:
            sorties = convertisseur(fichier, dossier_sortie)
            sorties = [partie for s in sorties
                       for partie in decouper_fichier(s, max_mots)]
            noms_sorties = "; ".join(s.name for s in sorties)
            print(f"  [OK     ] {relatif} -> {noms_sorties}")
            manifeste.append({
                "fichier_source": str(relatif), "categorie": categorie,
                "statut": "converti", "fichiers_sortie": noms_sorties, "erreur": "",
            })
            nb_ok += 1
        except Exception as exc:  # un fichier corrompu ne doit pas arrêter le lot
            print(f"  [ERREUR ] {relatif} : {exc}")
            manifeste.append({
                "fichier_source": str(relatif), "categorie": categorie,
                "statut": "erreur", "fichiers_sortie": "", "erreur": str(exc),
            })
            nb_erreurs += 1

    chemin_manifeste = dossier_sortie / "_MANIFEST.csv"
    with chemin_manifeste.open("w", encoding="utf-8", newline="") as f:
        champs = ["fichier_source", "categorie", "statut", "fichiers_sortie", "erreur"]
        redacteur = csv.DictWriter(f, fieldnames=champs, lineterminator="\n")
        redacteur.writeheader()
        redacteur.writerows(manifeste)

    print(f"\nBilan : {nb_ok} converti(s), {nb_ignores} ignoré(s), {nb_erreurs} erreur(s).")
    print(f"Manifeste : {chemin_manifeste}")
    return 0 if nb_erreurs == 0 else 2


def main() -> int:
    parseur = argparse.ArgumentParser(
        description="Normalise un dossier de fichiers hétérogènes en TXT/CSV UTF-8.")
    parseur.add_argument("--entree", default="sources_brutes",
                         help="Dossier d'entrée (défaut : sources_brutes)")
    parseur.add_argument("--sortie", default="sources_normalisees",
                         help="Dossier de sortie (défaut : sources_normalisees)")
    parseur.add_argument("--max-mots", type=int, default=MAX_MOTS_DEFAUT,
                         help=f"Nombre maximal de mots par fichier de sortie avant "
                              f"découpage (défaut : {MAX_MOTS_DEFAUT} ; 0 = désactivé)")
    arguments = parseur.parse_args()
    return normaliser_dossier(Path(arguments.entree), Path(arguments.sortie),
                              arguments.max_mots)


if __name__ == "__main__":
    sys.exit(main())
